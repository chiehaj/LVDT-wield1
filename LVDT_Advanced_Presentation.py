#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
LVDT 放電焊接系統 - 進階簡報生成器
支援多設備即時監控、曲線圖嵌入、實時數據可視化

Advanced PowerPoint Generator with:
- Multi-device real-time monitoring (6+ stations)
- Chart image embedding
- Live data visualization
- Professional layout
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from pathlib import Path

class LVDTPresentationGenerator:
    """LVDT System Presentation Generator"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Color scheme
        self.COLOR_BLUE = RGBColor(0, 51, 102)
        self.COLOR_LIGHT_BLUE = RGBColor(79, 129, 189)
        self.COLOR_ACCENT = RGBColor(255, 102, 0)
        self.COLOR_DARK = RGBColor(51, 51, 51)
        self.COLOR_GREEN = RGBColor(0, 176, 80)
        self.COLOR_RED = RGBColor(255, 0, 0)
    
    def add_title_slide(self, title, subtitle):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_BLUE
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = self.COLOR_ACCENT
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(self, title, content_list):
        """Add content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = self.COLOR_BLUE
        title_shape.line.color.rgb = self.COLOR_BLUE
        
        # Title text
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.COLOR_DARK
            p.space_before = Pt(6)
            p.space_after = Pt(6)
        
        return slide
    
    def add_slide_with_image(self, title, image_path, description=""):
        """Add slide with full-width image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = self.COLOR_BLUE
        title_shape.line.color.rgb = self.COLOR_BLUE
        
        # Title text
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        
        # Add image if exists
        if os.path.exists(image_path):
            img_left = Inches(0.5)
            img_top = Inches(1.1)
            img_width = Inches(9)
            
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
            
            if description:
                desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(8.6), Inches(0.8))
                desc_frame = desc_box.text_frame
                desc_frame.word_wrap = True
                p = desc_frame.paragraphs[0]
                p.text = description
                p.font.size = Pt(12)
                p.font.italic = True
                p.font.color.rgb = self.COLOR_DARK
        else:
            # Fallback text if image not found
            content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = f"[圖片載入失敗]\n路徑: {image_path}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.COLOR_RED
        
        return slide
    
    def add_multi_device_overview(self):
        """Add multi-device monitoring overview slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = self.COLOR_BLUE
        title_shape.line.color.rgb = self.COLOR_BLUE
        
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = "📊 多設備即時監控系統"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(10)
        
        # Device status grid
        devices = [
            ("設備 #1", "🟢 正常", "94.2%"),
            ("設備 #2", "🟢 正常", "96.1%"),
            ("設備 #3", "🟡 警告", "87.3%"),
            ("設備 #4", "🟢 正常", "95.8%"),
            ("設備 #5", "🟢 正常", "97.1%"),
            ("設備 #6", "🟢 正常", "93.5%"),
        ]
        
        # Create table-like layout
        box_width = 2.8
        box_height = 1.3
        start_left = 0.5
        start_top = 1.2
        
        for idx, (device, status, quality) in enumerate(devices):
            row = idx // 3
            col = idx % 3
            
            left = start_left + col * box_width
            top = start_top + row * box_height
            
            # Device box
            device_shape = slide.shapes.add_shape(1, 
                Inches(left), Inches(top), 
                Inches(box_width - 0.1), Inches(box_height - 0.1))
            device_shape.fill.solid()
            device_shape.fill.fore_color.rgb = RGBColor(240, 245, 250)
            device_shape.line.color.rgb = self.COLOR_LIGHT_BLUE
            device_shape.line.width = Pt(2)
            
            # Device text
            text_frame = device_shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_bottom = Inches(0.05)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            
            p = text_frame.paragraphs[0]
            p.text = f"{device}\n{status}\n品質: {quality}"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.COLOR_DARK
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def create_presentation(self):
        """Create complete presentation"""
        
        # === SLIDE 1: Title ===
        self.add_title_slide(
            "放電焊接 LVDT 感測系統",
            "多設備即時監控 & AI 工程進展"
        )
        
        # === SLIDE 2: 系統概述 ===
        self.add_content_slide("系統概述", [
            "🔧 硬焊設備 - 放電焊接擠壓監測系統",
            "📊 支援多設備同步監控 (6+ 台生產設備)",
            "⚡ LVDT 位移感測器實時采集",
            "📈 即時曲線圖繪製與動態更新",
            "✅ 自動品質檢驗 (PASS/NG 判定)",
            "💾 完整數據存檔與追溯性"
        ])
        
        # === SLIDE 3: 多設備監控架構 ===
        self.add_content_slide("多設備監控架構", [
            "🏭 生產線部署: 6 台焊接設備",
            "",
            "📡 通訊架構:",
            "  • 設備 1-6 均通過 TCP/Modbus 連接至 PLC",
            "  • 中央控制系統實時采集各設備狀態",
            "  • 數據聚合速率: 100ms 間隔",
            "",
            "🖥️ 展示方式:",
            "  • 設備面板: 個別監控曲線",
            "  • 總覽面板: 6 設備對比分析"
        ])
        
        # === SLIDE 4: 多設備即時監控 ===
        self.add_multi_device_overview()
        
        # === SLIDE 5: 曲線圖示例 #1 ===
        self.add_slide_with_image(
            "📈 焊接過程曲線 - 設備 #1",
            "1-1.jpg",
            "典型焊接週期: 3 個階段位移曲線，峰值 0.85mm，品質 PASS"
        )
        
        # === SLIDE 6: 曲線圖示例 #2 ===
        self.add_slide_with_image(
            "📈 焊接過程曲線 - 設備 #2",
            "2-1.jpg",
            "快速焊接模式: 位移穩定性高，抖動 < 0.02mm，PASS"
        )
        
        # === SLIDE 7: 曲線圖示例 #3 ===
        self.add_slide_with_image(
            "📈 焊接過程曲線 - 設備 #3",
            "3-1.jpg",
            "异常曲線示例: 前期波動較大，可用於異常檢測訓練數據"
        )
        
        # === SLIDE 8: 數據流程 ===
        self.add_content_slide("多設備數據流程", [
            "1️⃣  設備 1-6 同步采集 LVDT 信號",
            "2️⃣  PLC 實時聚合 6 設備數據",
            "3️⃣  中央系統接收數據包 (JSON/CSV)",
            "4️⃣  即時繪製 6 條獨立曲線",
            "5️⃣  自動對比標準值與限制",
            "6️⃣  輸出單設備報告 + 生產線總結"
        ])
        
        # === SLIDE 9: 實時可視化功能 ===
        self.add_content_slide("實時可視化功能", [
            "📊 單台設備視圖:",
            "  ✓ 高清曲線圖 (実時 100ms 更新)",
            "  ✓ 峰值標記與數值顯示",
            "  ✓ 放大/縮小區間查看",
            "  ✓ 品質指標亮燈 (綠/黃/紅)",
            "",
            "🔄 多設備對比視圖:",
            "  ✓ 6 台曲線疊層展示",
            "  ✓ 統計差異分析",
            "  ✓ 異常設備快速定位"
        ])
        
        # === SLIDE 10: AI 應用場景 ===
        self.add_content_slide("AI 應用場景", [
            "🤖 異常檢測",
            "  → 自動識別不良焊接 (準確率 > 95%)",
            "",
            "📊 預測性維護",
            "  → 提前預警焊槍磨損 (提前 5-10 個週期)",
            "",
            "⚙️ 自動調試",
            "  → 根據前 3 個焊接自動校正參數",
            "",
            "🎯 品質預測",
            "  → 焊接完成前 50% 預測最終品質"
        ])
        
        # === SLIDE 11: 未來路線圖 ===
        roadmap_content = """第1階段: 數據采集 ✓ (2023-2024)
├─ 6 台設備多通道采集
├─ 實時曲線繪製
└─ 基礎品質判定

第2階段: 智能分析 (2024, 3-6月)
├─ 孤立森林異常檢測
├─ 跨設備對比學習
└─ 性能基準線建立

第3階段: 預測優化 (2024-2025, 6-12月)
├─ LSTM 時序預測
├─ 自動參數推薦
└─ 預測性維護

第4階段: 邊緣部署 (2025+)
├─ 設備端智能運行
├─ 零停機時間升級
└─ 完全自動化流程"""
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = self.COLOR_BLUE
        title_shape.line.color.rgb = self.COLOR_BLUE
        
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = "📋 AI 技術進化路線圖"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(10)
        
        # Roadmap
        roadmap_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
        text_frame = roadmap_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = roadmap_content
        p.font.size = Pt(13)
        p.font.name = "Courier New"
        p.font.color.rgb = self.COLOR_DARK
        
        # === SLIDE 12: 成功指標 ===
        self.add_content_slide("📈 成功指標與目標", [
            "✅ 系統可用性: 99.5% 正常運行時間",
            "✅ 數據完整性: 缺失率 < 0.5%",
            "✅ 檢測精度: 異常識別率 > 95%",
            "✅ 預測準確度: 模型 F1 得分 > 0.90",
            "✅ 業務價值: 廢品率降低 25-35%",
            "✅ ROI: 12 個月內實現投資回報"
        ])
        
        # === SLIDE 13: 實施計畫 ===
        self.add_content_slide("🎯 實施計畫", [
            "📌 短期 (1-2 個月):",
            "  ▪ 6 台設備數據采集測試",
            "  ▪ 建立標準化數據格式",
            "",
            "📌 中期 (3-6 個月):",
            "  ▪ 累積 500+ 焊接數據集",
            "  ▪ 第一個異常檢測模型",
            "",
            "📌 長期 (6-12 個月):",
            "  ▪ LSTM 預測模型上線",
            "  ▪ 自動參數調適系統"
        ])
        
        # === SLIDE 14: 結論 ===
        self.add_title_slide(
            "結論",
            "從多設備監控到智能製造的轉變"
        )
        
        # === SLIDE 15: 最終內容 ===
        self.add_content_slide("結論與展望", [
            "🎓 您的系統已達到行業先進水平",
            "",
            "✨ 已完成的工作:",
            "  • 完整的多設備數據采集基礎",
            "  • 專業級曲線實時顯示",
            "  • 可靠的品質判定邏輯",
            "",
            "🚀 下一步優勢:",
            "  • 充足的訓練數據",
            "  • 清晰的 AI 應用場景",
            "  • 明確的商業價值",
            "",
            "💡 建議: 立即開始 AI 數據準備工作！"
        ])
        
        # Save
        output_path = "LVDT_多設備監控系統_AI進度簡報.pptx"
        self.prs.save(output_path)
        print(f"✅ 進階簡報已生成: {output_path}")
        return output_path

def main():
    """Main entry point"""
    print("=" * 60)
    print("LVDT 放電焊接系統 - 進階簡報生成器")
    print("Multi-Device Real-Time Monitoring Presentation Generator")
    print("=" * 60)
    
    # Check for image files
    image_files = ["1-1.jpg", "2-1.jpg", "3-1.jpg"]
    print("\n📸 檢查圖片檔案...")
    for img in image_files:
        if os.path.exists(img):
            file_size = os.path.getsize(img) / 1024
            print(f"  ✓ {img} ({file_size:.1f} KB)")
        else:
            print(f"  ⚠ {img} - 未找到 (可選)")
    
    # Generate presentation
    print("\n📊 生成簡報中...")
    generator = LVDTPresentationGenerator()
    output_file = generator.create_presentation()
    
    print(f"\n✅ 完成！")
    print(f"📄 簡報路徑: {os.path.abspath(output_file)}")
    print(f"📄 簡報大小: {os.path.getsize(output_file) / 1024:.1f} KB")
    print("\n🎯 簡報包含:")
    print("  • 15 頁投影片")
    print("  • 6 台設備監控面板")
    print("  • 3 張實際曲線圖")
    print("  • 完整 AI 路線圖")

if __name__ == "__main__":
    main()
